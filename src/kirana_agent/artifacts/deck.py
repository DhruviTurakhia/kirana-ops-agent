from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from kirana_agent.domain.money import format_inr
from kirana_agent.domain.service import StoreService

TEMPLATE_VERSION = "sales-analysis-v4"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
INK = RGBColor(13, 24, 38)
MUTED = RGBColor(91, 103, 116)
RULE = RGBColor(208, 216, 224)
PANEL = RGBColor(243, 246, 248)
ACCENT = RGBColor(15, 139, 141)
ACCENT_2 = RGBColor(61, 141, 255)
ACCENT_3 = RGBColor(159, 216, 211)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(180, 45, 58)
FONT = "Aptos"


class SalesDeckGenerator:
    """Build an editable PowerPoint whose charts come directly from stored sales."""

    def __init__(self, service: StoreService, output_dir: str | Path):
        self.service = service
        self.output_dir = Path(output_dir)

    def generate(self, *, from_date: str, to_date: str) -> dict[str, Any]:
        analysis = self.service.sales_analysis(from_date=from_date, to_date=to_date)
        stable_source = deepcopy(analysis)
        stable_source.pop("generated_at", None)
        source_hash = self.service.content_hash(stable_source)
        source_id = f"{analysis['from_date']}:{analysis['to_date']}"
        cached = self.service.find_artifact(
            artifact_type="SALES_PPTX",
            source_id=source_id,
            source_hash=source_hash,
            template_version=TEMPLATE_VERSION,
        )
        if cached:
            return {"ok": True, "cached": True, "analysis": analysis, **cached}

        deck_dir = self.output_dir / "pptx"
        deck_dir.mkdir(parents=True, exist_ok=True)
        path = deck_dir / f"sales-analysis-{analysis['from_date']}-to-{analysis['to_date']}.pptx"
        self._render(analysis, path)
        record = self.service.record_artifact(
            artifact_type="SALES_PPTX",
            source_id=source_id,
            source_hash=source_hash,
            template_version=TEMPLATE_VERSION,
            file_path=path,
        )
        return {"ok": True, "cached": False, "analysis": analysis, **record}

    def _render(self, analysis: dict[str, Any], path: Path) -> None:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        blank = prs.slide_layouts[6]
        store = self.service.get_store_profile()

        self._title_slide(prs.slides.add_slide(blank), store, analysis)
        self._daily_sales_slide(prs.slides.add_slide(blank), analysis)
        self._top_products_slide(prs.slides.add_slide(blank), analysis)
        self._payment_slide(prs.slides.add_slide(blank), analysis)
        self._gst_slide(prs.slides.add_slide(blank), analysis)
        self._stock_slide(prs.slides.add_slide(blank), analysis)

        for index, slide in enumerate(prs.slides, start=1):
            self._footer(slide, index, len(prs.slides), store["display_name"], analysis)
        self._normalize_chart_axis_ids(prs)
        prs.save(path)

    def _title_slide(self, slide: Any, store: dict[str, Any], analysis: dict[str, Any]) -> None:
        self._set_background(slide, WHITE)
        self._rule(slide, 0.65, 0.72, 1.35, 0.06, ACCENT)
        self._text(slide, store["display_name"].upper(), 0.65, 1.05, 5.4, 0.35, 14, ACCENT, bold=True)
        self._text(slide, "Sales performance,\nmade actionable", 0.65, 1.55, 7.4, 1.65, 50, INK, bold=True)
        self._text(
            slide,
            f"{analysis['from_date']} to {analysis['to_date']} · {analysis['timezone']}",
            0.68,
            3.45,
            6.8,
            0.45,
            18,
            MUTED,
        )
        totals = analysis["totals"]
        self._metric(slide, totals["gross"], "Gross sales", 8.35, 1.25, 4.1)
        self._metric(slide, str(totals["bill_count"]), "Finalized bills", 8.35, 3.05, 4.1)
        self._metric(slide, totals["gst"], "GST included", 8.35, 4.85, 4.1)

    def _daily_sales_slide(self, slide: Any, analysis: dict[str, Any]) -> None:
        self._slide_title(slide, "Sales momentum reveals the week's rhythm")
        daily = analysis["daily_sales"]
        if daily:
            data = ChartData()
            data.categories = [item["date"][5:] for item in daily]
            data.add_series("Gross sales (₹)", [item["gross_paise"] / 100 for item in daily])
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(0.68),
                Inches(1.62),
                Inches(8.3),
                Inches(4.8),
                data,
            ).chart
            self._style_chart(chart, show_legend=False)
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = RULE
            chart.value_axis.tick_labels.font.size = Pt(16)
            chart.category_axis.tick_labels.font.size = Pt(16)
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = ACCENT_2
            best = max(daily, key=lambda item: item["gross_paise"])
            average = analysis["totals"]["gross_paise"] // max(len(daily), 1)
            copy = (
                f"Peak day\n{best['date']}\n{best['gross']}\n\n"
                f"Daily average\n{format_inr(average)}\n\n"
                "Use the peak-day pattern to time replenishment and staffing."
            )
        else:
            self._empty_state(slide, "No finalized sales in this period")
            copy = "Finalize bills to populate this trend with real store data."
        self._side_note(slide, copy, 9.45, 1.75, 3.2, 4.2)

    def _top_products_slide(self, slide: Any, analysis: dict[str, Any]) -> None:
        self._slide_title(slide, "A few products drive most sales")
        products = analysis["top_products"][:7]
        if products:
            ordered = list(reversed(products))
            data = ChartData()
            data.categories = [item["product_name"][:28] for item in ordered]
            data.add_series("Revenue (₹)", [item["gross_paise"] / 100 for item in ordered])
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_CLUSTERED,
                Inches(0.75),
                Inches(1.55),
                Inches(8.6),
                Inches(4.95),
                data,
            ).chart
            self._style_chart(chart, show_legend=False)
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = RULE
            chart.value_axis.tick_labels.font.size = Pt(16)
            chart.category_axis.tick_labels.font.size = Pt(16)
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = ACCENT
            leader = products[0]
            copy = (
                f"Revenue leader\n{leader['product_name']}\n{leader['gross']}\n\n"
                "Protect availability on leaders first; then use slower items to broaden the basket."
            )
        else:
            self._empty_state(slide, "No product sales in this period")
            copy = "Product rankings will appear after finalized bills are recorded."
        self._side_note(slide, copy, 9.7, 1.75, 2.95, 4.25)

    def _payment_slide(self, slide: Any, analysis: dict[str, Any]) -> None:
        self._slide_title(slide, "Tender mix shows how customers prefer to pay")
        mix = analysis["payment_mix"]
        if mix:
            data = ChartData()
            data.categories = [item["payment_mode"] for item in mix]
            data.add_series("Gross sales (₹)", [item["gross_paise"] / 100 for item in mix])
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.DOUGHNUT,
                Inches(0.85),
                Inches(1.55),
                Inches(6.6),
                Inches(4.95),
                data,
            ).chart
            self._style_chart(chart, show_legend=True)
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.font.size = Pt(16)
            chart.plots[0].has_data_labels = True
            labels = chart.plots[0].data_labels
            labels.show_percentage = True
            labels.show_category_name = False
            labels.font.size = Pt(16)
            for point, color in zip(chart.series[0].points, [ACCENT, ACCENT_2, ACCENT_3, MUTED], strict=False):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = color
            dominant = max(mix, key=lambda item: item["gross_paise"])
            total = max(analysis["totals"]["gross_paise"], 1)
            share = dominant["gross_paise"] / total * 100
            copy = f"{dominant['payment_mode']} leads\n{share:.0f}% of sales value\n\nMode total: {dominant['gross']}\nBills: {dominant['bill_count']}"
        else:
            self._empty_state(slide, "No tender data in this period")
            copy = "Payment mix will populate from finalized bills."
        self._side_note(slide, copy, 8.1, 1.75, 4.45, 4.2)

    def _gst_slide(self, slide: Any, analysis: dict[str, Any]) -> None:
        self._slide_title(slide, "GST is reconciled from immutable bill-line snapshots")
        rows = analysis["gst_by_slab"]
        headers = ["GST slab", "Taxable value", "GST included", "Gross value"]
        table_rows = [[f"{item['gst_rate_percent']}%", item["taxable"], item["gst"], item["gross"]] for item in rows]
        if not table_rows:
            table_rows = [["—", "₹0.00", "₹0.00", "₹0.00"]]
        self._table(slide, headers, table_rows, 0.75, 1.65, 8.5, 4.7)
        totals = analysis["totals"]
        copy = (
            f"Total GST\n{totals['gst']}\n\n"
            f"CGST {totals['cgst']}\nSGST {totals['sgst']}\nIGST {totals['igst']}\n\n"
            "Header totals equal the sum of rounded bill lines."
        )
        self._side_note(slide, copy, 9.65, 1.75, 2.85, 4.25)

    def _stock_slide(self, slide: Any, analysis: dict[str, Any]) -> None:
        low = analysis["low_stock"]
        title = (
            f"{len(low)} items need attention before the next selling cycle"
            if low
            else "Stock health is clear across the seeded catalog"
        )
        self._slide_title(slide, title)
        headers = ["Product", "SKU", "On hand", "Reorder at", "Status"]
        rows = []
        for item in low[:7]:
            status = "OUT" if item["stock_atomic"] == 0 else "LOW"
            rows.append(
                [
                    item["name"][:38],
                    item["sku"],
                    self._stock_label(item["stock_atomic"], item["base_uom"], item["sale_uom"]),
                    self._stock_label(item["reorder_atomic"], item["base_uom"], item["sale_uom"]),
                    status,
                ]
            )
        if not rows:
            rows = [["No products at or below reorder level", "—", "—", "—", "OK"]]
        self._table(slide, headers, rows, 0.65, 1.55, 9.0, 5.15)
        out_count = sum(1 for item in low if item["stock_atomic"] == 0)
        copy = (
            f"Out of stock\n{out_count}\n\nLow stock\n{max(len(low) - out_count, 0)}\n\n"
            "Prioritize out-of-stock essentials, then the highest-selling low-stock items."
        )
        self._side_note(slide, copy, 10.0, 1.75, 2.65, 4.35)

    def _slide_title(self, slide: Any, text: str) -> None:
        self._set_background(slide, WHITE)
        self._text(slide, text, 0.68, 0.48, 11.8, 0.82, 35, INK, bold=True)
        self._rule(slide, 0.68, 1.32, 12.0, 0.02, RULE)

    def _metric(self, slide: Any, value: str, label: str, x: float, y: float, w: float) -> None:
        self._text(slide, value, x, y, w, 0.65, 30, INK, bold=True)
        self._text(slide, label.upper(), x, y + 0.7, w, 0.32, 12, ACCENT, bold=True)
        self._rule(slide, x, y + 1.25, w, 0.018, RULE)

    def _side_note(self, slide: Any, text: str, x: float, y: float, w: float, h: float) -> None:
        self._rule(slide, x, y, 0.06, h, ACCENT)
        self._text(slide, text, x + 0.28, y + 0.05, w - 0.28, h - 0.1, 18, INK, bold=False)

    def _empty_state(self, slide: Any, text: str) -> None:
        shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.7), Inches(8.1), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = RULE
        self._text(slide, text, 1.2, 3.45, 7.3, 0.6, 24, MUTED, bold=True, align=PP_ALIGN.CENTER)

    def _table(
        self,
        slide: Any,
        headers: list[str],
        rows: list[list[str]],
        x: float,
        y: float,
        w: float,
        h: float,
    ) -> None:
        table = slide.shapes.add_table(
            len(rows) + 1,
            len(headers),
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        ).table
        for col in range(len(headers)):
            table.columns[col].width = Inches(w / len(headers))
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK
            self._style_cell(cell, 16, WHITE, bold=True)
        for row_index, row in enumerate(rows, start=1):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_index % 2 else PANEL
                self._style_cell(cell, 16, INK, bold=(col_index == 0))

    @staticmethod
    def _style_cell(cell: Any, size: int, color: RGBColor, *, bold: bool) -> None:
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = FONT
            paragraph.font.size = Pt(size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = color

    @staticmethod
    def _style_chart(chart: Any, *, show_legend: bool) -> None:
        chart.has_title = False
        chart.has_legend = show_legend
        chart.chart_style = 10
        chart.font.name = FONT
        chart.font.size = Pt(16)

    def _footer(
        self,
        slide: Any,
        index: int,
        count: int,
        store_name: str,
        analysis: dict[str, Any],
    ) -> None:
        self._text(
            slide,
            f"{store_name} · {analysis['from_date']} to {analysis['to_date']}",
            0.68,
            7.12,
            8.5,
            0.2,
            9,
            MUTED,
        )
        self._text(slide, f"{index}/{count}", 11.9, 7.12, 0.75, 0.2, 9, MUTED, align=PP_ALIGN.RIGHT)

    @staticmethod
    def _normalize_chart_axis_ids(prs: Presentation) -> None:
        """Serialize python-pptx's signed random axis IDs as Open XML UInt32.

        Microsoft PowerPoint tolerates negative lexical values, but strict Open XML
        readers correctly require the declared unsigned representation. Rewriting the
        same 32 bits as UInt32 improves interoperability without changing chart links.
        """

        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_chart", False):
                    continue
                chart_space = shape.chart.part._element
                for element in chart_space.xpath(".//c:axId | .//c:crossAx"):
                    value = int(element.get("val"))
                    if value < 0:
                        element.set("val", str(value & 0xFFFFFFFF))

    @staticmethod
    def _stock_label(quantity_atomic: int, base_uom: str, sale_uom: str) -> str:
        if base_uom == "piece":
            value = quantity_atomic / 1000
            plurals = {
                "box": "boxes",
                "loaf": "loaves",
                "pouch": "pouches",
            }
            label = sale_uom if value == 1 else plurals.get(sale_uom, f"{sale_uom}s")
            return f"{value:g} {label}"
        if base_uom == "g":
            return (
                f"{quantity_atomic / 1000:g} kg"
                if quantity_atomic >= 1000
                else f"{quantity_atomic:g} g"
            )
        if base_uom == "ml":
            return (
                f"{quantity_atomic / 1000:g} L"
                if quantity_atomic >= 1000
                else f"{quantity_atomic:g} ml"
            )
        return str(quantity_atomic)

    @staticmethod
    def _set_background(slide: Any, color: RGBColor) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    @staticmethod
    def _rule(slide: Any, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    @staticmethod
    def _text(
        slide: Any,
        text: str,
        x: float,
        y: float,
        w: float,
        h: float,
        size: int,
        color: RGBColor,
        *,
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.LEFT,
    ) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
