from __future__ import annotations

from datetime import datetime
from pathlib import Path
from zipfile import ZipFile

import pytest

pytest.importorskip("pptx")

from pptx import Presentation

from kirana_agent.artifacts.deck import SalesDeckGenerator


def test_sales_deck_is_cached_openable_and_contains_real_chart_data(store, tmp_path) -> None:
    first_product = store.create_stocked_product(
        gst_rate_bps=500, sell_price="105.00", label="deck-five"
    )
    second_product = store.create_stocked_product(
        gst_rate_bps=1800, sell_price="118.00", label="deck-eighteen"
    )
    cash_bill = store.finalize_sale(
        product_id=first_product["id"], quantity=2, chat_id="deck-cash"
    )
    upi_bill = store.finalize_sale(
        product_id=second_product["id"],
        quantity=1,
        chat_id="deck-upi",
        payment_mode="UPI",
        payment_reference="UPI-DECK",
    )
    today = datetime.now(store.service.timezone).date().isoformat()
    generator = SalesDeckGenerator(store.service, tmp_path / "artifacts")

    first = generator.generate(from_date=today, to_date=today)
    second = generator.generate(from_date=today, to_date=today)

    assert first["cached"] is False
    assert second["cached"] is True
    assert second["id"] == first["id"]
    path = Path(first["file_path"])
    assert path.is_file()
    with ZipFile(path) as archive:
        assert archive.testzip() is None
        names = set(archive.namelist())
        assert "[Content_Types].xml" in names
        assert "ppt/presentation.xml" in names
        assert len([name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")]) == 6

    presentation = Presentation(path)
    assert len(presentation.slides) == 6
    charts = [
        shape.chart
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_chart", False)
    ]
    assert len(charts) >= 3
    value_vectors = [tuple(round(float(value), 2) for value in chart.series[0].values) for chart in charts]
    expected_sales = sorted([cash_bill["gross_paise"] / 100, upi_bill["gross_paise"] / 100])
    assert any(sorted(vector) == expected_sales for vector in value_vectors)
    assert any(vector == (sum(expected_sales),) for vector in value_vectors)
    assert first["analysis"]["totals"]["bill_count"] == 2
    assert first["analysis"]["totals"]["gross_paise"] == cash_bill["gross_paise"] + upi_bill["gross_paise"]
    assert store.scalar(
        "SELECT COUNT(*) FROM artifacts WHERE artifact_type = 'SALES_PPTX'"
    ) == 1
